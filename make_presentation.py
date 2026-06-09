"""
make_presentation.py
Generates the final project presentation (PPTX) for the Touchless Writing System
— Image Processing & Computer Vision course project.

Run:  python make_presentation.py
Out:  Touchless_Writing_System.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Theme
# ----------------------------------------------------------------------------
INK     = RGBColor(0x0F, 0x1B, 0x2A)   # deep navy (primary dark)
TEAL    = RGBColor(0x14, 0xB8, 0xA6)   # accent teal
TEAL_D  = RGBColor(0x0D, 0x7A, 0x6E)   # darker teal
AMBER   = RGBColor(0xF5, 0xA6, 0x23)   # accent amber
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MIST    = RGBColor(0xF4, 0xF7, 0xF9)   # light panel
SLATE   = RGBColor(0x3C, 0x4A, 0x5A)   # body text grey
LINE    = RGBColor(0xD7, 0xDF, 0xE6)

FONT_H  = "Calibri"
FONT_B  = "Calibri"

SW, SH = Inches(13.333), Inches(7.5)   # 16:9

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# helpers
# ----------------------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of (text, size, bold, color, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, b, c, *rest) in para:
            it = rest[0] if rest else False
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.bold = b
            r.font.italic = it
            r.font.name = FONT_B
            r.font.color.rgb = c
    return tb


def chip(s, x, y, label, fill=TEAL, fg=WHITE, w=None, sz=12):
    w = w or Inches(0.05 + 0.105 * len(label))
    c = rect(s, x, y, w, Inches(0.34), fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    c.adjustments[0] = 0.5
    tf = c.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(sz)
    r.font.bold = True
    r.font.color.rgb = fg
    r.font.name = FONT_B
    return c


def header(s, kicker, title, num):
    """Standard content-slide header band."""
    rect(s, 0, 0, SW, Inches(1.35), INK)
    rect(s, 0, Inches(1.35), SW, Inches(0.06), TEAL)
    rect(s, Inches(0.55), Inches(0.34), Inches(0.12), Inches(0.66), AMBER)
    txt(s, Inches(0.85), Inches(0.30), Inches(10.5), Inches(0.3),
        [[(kicker.upper(), 12.5, True, TEAL)]])
    txt(s, Inches(0.85), Inches(0.58), Inches(11.4), Inches(0.7),
        [[(title, 27, True, WHITE)]])
    # slide number badge
    b = rect(s, Inches(12.45), Inches(0.46), Inches(0.42), Inches(0.42),
             TEAL, shape=MSO_SHAPE.OVAL)
    tf = b.text_frame
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num)
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE


def bullets(s, x, y, w, items, sz=15, gap=10, color=SLATE, marker=TEAL):
    tb = s.shapes.add_textbox(x, y, w, Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0
    for i, item in enumerate(items):
        # item: (text) or (lead, rest) tuple
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        rb = p.add_run()
        rb.text = "▸  "
        rb.font.size = Pt(sz)
        rb.font.bold = True
        rb.font.color.rgb = marker
        rb.font.name = FONT_B
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run(); r1.text = lead
            r1.font.size = Pt(sz); r1.font.bold = True
            r1.font.color.rgb = INK; r1.font.name = FONT_B
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(sz); r2.font.color.rgb = color
            r2.font.name = FONT_B
        else:
            r = p.add_run(); r.text = item
            r.font.size = Pt(sz); r.font.color.rgb = color
            r.font.name = FONT_B
    return tb


def card(s, x, y, w, h, title, body, accent=TEAL, icon=""):
    panel = rect(s, x, y, w, h, WHITE, line=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    panel.adjustments[0] = 0.06
    rect(s, x, y, Inches(0.14), h, accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    head = (icon + "  " + title) if icon else title
    txt(s, x + Inches(0.32), y + Inches(0.18), w - Inches(0.5), Inches(0.4),
        [[(head, 15.5, True, INK)]])
    txt(s, x + Inches(0.32), y + Inches(0.62), w - Inches(0.55), h - Inches(0.7),
        [[(body, 12.5, False, SLATE)]], line_spacing=1.08)
    return panel


# ============================================================================
# SLIDE 1 — COVER
# ============================================================================
s = slide()
bg(s, INK)
# decorative shapes
rect(s, Inches(9.3), 0, Inches(4.05), SH, RGBColor(0x12, 0x24, 0x38))
rect(s, Inches(9.3), 0, Inches(0.06), SH, TEAL)
# big translucent rings (concentric) on right
for d, col in [(5.2, RGBColor(0x17, 0x30, 0x46)), (3.7, RGBColor(0x19, 0x39, 0x52)),
               (2.2, TEAL_D)]:
    o = rect(s, Inches(11.35 - d/2), Inches(3.75 - d/2), Inches(d), Inches(d),
             col, shape=MSO_SHAPE.OVAL)
txt(s, Inches(10.05), Inches(3.05), Inches(2.7), Inches(1.5),
    [[("✍︎", 46, True, WHITE)], [("AIR", 17, True, TEAL)],
     [("CANVAS", 17, True, WHITE)]], align=PP_ALIGN.CENTER)

chip(s, Inches(0.85), Inches(0.95), "IMAGE PROCESSING & COMPUTER VISION", fill=TEAL, sz=12, w=Inches(4.0))
txt(s, Inches(0.85), Inches(1.9), Inches(8.2), Inches(2.4),
    [[("Touchless", 54, True, WHITE)],
     [("Writing System", 54, True, TEAL)]], line_spacing=0.98)
txt(s, Inches(0.88), Inches(4.05), Inches(8.0), Inches(0.9),
    [[("Air-writing with real-time hand-gesture tracking,", 17, False, RGBColor(0xC7,0xD3,0xDE))],
     [("OCR handwriting recognition & voice dictation.", 17, False, RGBColor(0xC7,0xD3,0xDE))]],
    line_spacing=1.1)
rect(s, Inches(0.88), Inches(5.0), Inches(2.4), Inches(0.045), AMBER)

txt(s, Inches(0.88), Inches(5.35), Inches(8), Inches(1.3),
    [[("Presented by", 12, True, TEAL)],
     [("Mohammed Farees Mohammed Rifath", 18, True, WHITE)],
     [("Course Project  ·  Computer Vision", 13, False, RGBColor(0x9F,0xB0,0xBF))]],
    space_after=4)

# tech tags
tagx = Inches(0.88)
for tg in ["OpenCV", "MediaPipe", "NumPy", "OpenAI GPT-4o", "Python 3.11"]:
    c = chip(s, tagx, Inches(6.7), tg, fill=RGBColor(0x1C,0x32,0x49), fg=RGBColor(0xCF,0xDD,0xE8), sz=11)
    tagx = Emu(c.left + c.width + Inches(0.16))


# ============================================================================
# SLIDE 2 — AGENDA
# ============================================================================
s = slide(); bg(s, WHITE)
header(s, "Overview", "Agenda", 2)
items = [
    ("01", "Project Problem", "Why touch-free input matters", TEAL),
    ("02", "Literature Review", "Existing air-writing & OCR methods", AMBER),
    ("03", "Data Acquisition", "How input is captured", TEAL),
    ("04", "Methodology", "Pipeline: tracking → canvas → OCR", AMBER),
    ("05", "Results & Discussion", "CER / WER / latency evaluation", TEAL),
    ("06", "Novelty & Conclusion", "Contributions and future work", AMBER),
]
x0, y0 = Inches(0.85), Inches(1.85)
cw, ch = Inches(5.75), Inches(1.45)
gx, gy = Inches(0.55), Inches(0.35)
for i, (n, t, d, ac) in enumerate(items):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    p = rect(s, cx, cy, cw, ch, MIST, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    p.adjustments[0] = 0.08
    rect(s, cx, cy, Inches(0.14), ch, ac, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, cx + Inches(0.32), cy + Inches(0.2), Inches(1.5), Inches(1),
        [[(n, 34, True, ac)]])
    txt(s, cx + Inches(1.55), cy + Inches(0.26), cw - Inches(1.7), Inches(1.1),
        [[(t, 18, True, INK)], [(d, 12.5, False, SLATE)]], space_after=3)


# ============================================================================
# SLIDE 3 — PROJECT PROBLEM
# ============================================================================
s = slide(); bg(s, WHITE)
header(s, "01 · Motivation", "Project Problem", 3)
txt(s, Inches(0.85), Inches(1.7), Inches(11.6), Inches(0.7),
    [[("Conventional digital writing needs a ", 16, False, SLATE),
      ("physical touch surface", 16, True, INK),
      (" — keyboard, stylus or trackpad. That is limiting in many real settings.", 16, False, SLATE)]],
    line_spacing=1.1)

probs = [
    ("🖐️", "Hygiene & accessibility", "Shared touch devices spread germs; some users cannot operate a stylus or keyboard comfortably."),
    ("📊", "Hardware cost", "Digital pens, smart-boards and touch panels are expensive and not always available."),
    ("🎯", "Natural interaction", "There is no intuitive, hardware-light way to write/annotate in the air during teaching or presentations."),
    ("🔤", "Hand-off to digital", "Even when notes are captured, turning handwriting into editable, searchable text is a separate manual step."),
]
x0, y0 = Inches(0.85), Inches(2.55)
cw, ch = Inches(5.75), Inches(1.85)
for i, (ic, t, d) in enumerate(probs):
    cx = x0 + (cw + Inches(0.55)) * (i % 2)
    cy = y0 + (ch + Inches(0.3)) * (i // 2)
    card(s, cx, cy, cw, ch, t, d, accent=(TEAL if i % 2 == 0 else AMBER), icon=ic)

# problem statement strip
rect(s, Inches(0.85), Inches(6.75), Inches(11.6), Inches(0.5), INK,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.1), Inches(6.82), Inches(11.1), Inches(0.4),
    [[("Goal:  ", 14, True, TEAL),
      ("a webcam-only system to write in the air with hand gestures and convert that handwriting to digital text.",
       14, False, WHITE)]])


# ============================================================================
# SLIDE 4 — OBJECTIVES
# ============================================================================
s = slide(); bg(s, WHITE)
header(s, "01 · Scope", "Aim & Objectives", 4)
# left aim panel
p = rect(s, Inches(0.85), Inches(1.8), Inches(4.4), Inches(4.9), INK,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE); p.adjustments[0] = 0.05
rect(s, Inches(1.15), Inches(2.15), Inches(0.7), Inches(0.1), TEAL)
txt(s, Inches(1.15), Inches(2.4), Inches(3.85), Inches(0.5),
    [[("Project Aim", 20, True, WHITE)]])
txt(s, Inches(1.15), Inches(3.0), Inches(3.85), Inches(3.4),
    [[("Build a touchless writing application that lets a user ", 15, False, RGBColor(0xCF,0xDD,0xE8)),
      ("draw in the air", 15, True, TEAL),
      (" using only a standard webcam, and ", 15, False, RGBColor(0xCF,0xDD,0xE8)),
      ("convert the handwriting into editable digital text", 15, True, WHITE),
      (" using computer-vision pre-processing and OCR.", 15, False, RGBColor(0xCF,0xDD,0xE8))]],
    line_spacing=1.18)

objs = [
    ("Real-time hand tracking", "Detect 21 hand landmarks from webcam and recognise finger-state gestures."),
    ("Gesture-driven canvas", "Map gestures to move / draw / erase / select / clear on a digital canvas."),
    ("CV pre-processing pipeline", "Clean strokes (CLAHE, thresholding, morphology, crop) to aid recognition."),
    ("Handwriting → text (OCR)", "Convert the rendered page into editable, savable text."),
    ("Multimodal extras", "Voice-to-text dictation and live on-screen text overlay."),
    ("Objective evaluation", "Measure OCR accuracy (CER / WER) and latency, with vs. without pre-processing."),
]
x0, y0 = Inches(5.55), Inches(1.85)
for i, (t, d) in enumerate(objs):
    cy = y0 + (Inches(0.78) + Inches(0.04)) * i + Inches(0.0)
    cy = y0 + i * Inches(0.815)
    n = rect(s, x0, cy, Inches(0.42), Inches(0.42), TEAL, shape=MSO_SHAPE.OVAL)
    tf = n.text_frame; tf.margin_top=0; tf.margin_bottom=0
    pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    rr = pp.add_run(); rr.text = str(i+1); rr.font.bold=True; rr.font.size=Pt(15); rr.font.color.rgb=WHITE
    txt(s, x0 + Inches(0.62), cy - Inches(0.04), Inches(6.5), Inches(0.8),
        [[(t + " — ", 14.5, True, INK), (d, 13, False, SLATE)]], line_spacing=1.0)


# ============================================================================
# SLIDE 5 — LITERATURE REVIEW
# ============================================================================
s = slide(); bg(s, WHITE)
header(s, "02 · Background", "Literature Review — Existing Methodologies", 5)

cols = [
    ("Sensor / glove based", "Leap Motion, data-gloves, depth (Kinect) & wearable IMUs track the hand in 3-D.",
     "High accuracy", "Extra hardware, cost, not webcam-friendly", AMBER),
    ("Color-marker tracking", "Classic OpenCV air-canvas: track a colored cap/fingertip via HSV thresholding.",
     "Simple, real-time, cheap", "Needs a colored marker; fragile to lighting", TEAL),
    ("Deep-learning landmarks", "CNN hand-pose estimators — MediaPipe Hands, OpenPose — give 21 keypoints from RGB.",
     "Marker-free, robust, fast", "Needs a trained model / library", AMBER),
    ("OCR / HTR engines", "Tesseract (classical) and CNN-RNN / vision-LLM models read handwriting into text.",
     "Mature, off-the-shelf", "Tesseract weak on free handwriting", TEAL),
]
cw = Inches(2.92); ch = Inches(4.55); gap = Inches(0.13)
x0 = Inches(0.6); y0 = Inches(1.85)
for i, (t, d, pro, con, ac) in enumerate(cols):
    cx = x0 + (cw + gap) * i
    p = rect(s, cx, y0, cw, ch, MIST, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    p.adjustments[0] = 0.05
    rect(s, cx, y0, cw, Inches(0.85), ac, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, cx, y0 + Inches(0.55), cw, Inches(0.32), ac)  # square off bottom of header
    txt(s, cx + Inches(0.18), y0 + Inches(0.16), cw - Inches(0.3), Inches(0.7),
        [[(t, 14, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, cx + Inches(0.2), y0 + Inches(1.0), cw - Inches(0.36), Inches(1.6),
        [[(d, 11.8, False, SLATE)]], line_spacing=1.06)
    txt(s, cx + Inches(0.2), y0 + Inches(2.75), cw - Inches(0.36), Inches(0.9),
        [[("✓ ", 12, True, TEAL_D), (pro, 11.5, True, INK)]], line_spacing=1.0)
    txt(s, cx + Inches(0.2), y0 + Inches(3.55), cw - Inches(0.36), Inches(0.9),
        [[("✗ ", 12, True, RGBColor(0xC0,0x4A,0x3A)), (con, 11.5, False, SLATE)]], line_spacing=1.0)

txt(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.7),
    [[("Our positioning:  ", 14, True, AMBER),
      ("combine marker-free DL landmark tracking (MediaPipe) for input with a classical CV pre-processing pipeline + a vision-LLM OCR — webcam-only, no special hardware.",
       13.5, False, SLATE)]], line_spacing=1.08)


# ============================================================================
# SLIDE 6 — DATA ACQUISITION
# ============================================================================
s = slide(); bg(s, WHITE)
header(s, "03 · Input", "Data Acquisition", 6)

# left: capture flow
txt(s, Inches(0.85), Inches(1.65), Inches(6), Inches(0.4),
    [[("How input data enters the system", 15, True, INK)]])
flow = [
    ("Webcam stream", "Live RGB frames, 1280×720, captured with OpenCV (cv2.VideoCapture, index 0).", TEAL),
    ("Hand landmarks", "MediaPipe Hands extracts 21 (x, y, z) keypoints per frame — single hand.", AMBER),
    ("Gesture signal", "Finger up/down states → discrete gesture, stabilised over 3 frames.", TEAL),
    ("Canvas image", "Cursor trail rendered to a 1280×720 NumPy canvas — the 'page' data.", AMBER),
    ("Voice & evaluation", "Microphone audio (voice mode) + a small labelled image/text set for OCR eval.", TEAL),
]
y = Inches(2.15)
for i, (t, d, ac) in enumerate(flow):
    n = rect(s, Inches(0.85), y, Inches(0.46), Inches(0.46), ac, shape=MSO_SHAPE.OVAL)
    tf=n.text_frame; tf.margin_top=0; tf.margin_bottom=0
    pp=tf.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER
    rr=pp.add_run(); rr.text=str(i+1); rr.font.bold=True; rr.font.size=Pt(15); rr.font.color.rgb=WHITE
    if i < len(flow)-1:
        rect(s, Inches(1.06), y + Inches(0.46), Inches(0.04), Inches(0.42), LINE)
    txt(s, Inches(1.5), y - Inches(0.02), Inches(5.3), Inches(0.9),
        [[(t, 14.5, True, INK)], [(d, 12, False, SLATE)]], space_after=2, line_spacing=1.0)
    y = y + Inches(0.88)

# right: spec panel
px = Inches(7.5)
p = rect(s, px, Inches(1.65), Inches(5.0), Inches(5.1), INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
p.adjustments[0] = 0.05
txt(s, px + Inches(0.35), Inches(1.95), Inches(4.3), Inches(0.4),
    [[("Acquisition Specs", 17, True, TEAL)]])
specs = [
    ("Sensor", "Standard laptop / USB webcam — no depth, no glove"),
    ("Resolution", "1280 × 720 frames; mirrored for natural drawing"),
    ("Tracking model", "MediaPipe Hands (legacy Solutions API)"),
    ("Landmarks", "21 per hand · used to derive 6 gestures"),
    ("Canvas", "NumPy array, multi-page, 8 colors + eraser"),
    ("Audio", "Mic → Google Speech (voice dictation mode)"),
    ("Eval data", "5–10 saved page images + exact ground-truth .txt"),
]
yy = Inches(2.5)
for k, v in specs:
    txt(s, px + Inches(0.35), yy, Inches(1.55), Inches(0.6),
        [[(k, 12.5, True, AMBER)]])
    txt(s, px + Inches(1.85), yy, Inches(2.95), Inches(0.6),
        [[(v, 12, False, RGBColor(0xCF,0xDD,0xE8))]], line_spacing=1.0)
    yy += Inches(0.6)


# ============================================================================
# SLIDE 7 — METHODOLOGY: ARCHITECTURE
# ============================================================================
s = slide(); bg(s, WHITE)
header(s, "04 · Methodology", "System Architecture — End-to-End Pipeline", 7)

stages = [
    ("Webcam\nFrame", "RGB capture\n(OpenCV)", TEAL),
    ("Hand\nTracking", "21 landmarks\n(MediaPipe)", AMBER),
    ("Gesture\nDetection", "Finger states\n→ action", TEAL),
    ("Digital\nCanvas", "NumPy page\nrender", AMBER),
    ("CV Pre-\nprocess", "CLAHE · thresh\nmorph · crop", TEAL),
    ("OCR\nEngine", "GPT-4o-mini\n→ text", AMBER),
]
n = len(stages)
bw = Inches(1.78); bh = Inches(1.55); gap = Inches(0.27)
total = bw*n + gap*(n-1)
x0 = (SW - total) / 2
y0 = Inches(2.35)
for i, (t, d, ac) in enumerate(stages):
    cx = x0 + (bw + gap) * i
    box = rect(s, cx, y0, bw, bh, ac, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box.adjustments[0] = 0.09
    txt(s, cx + Inches(0.1), y0 + Inches(0.2), bw - Inches(0.2), Inches(0.8),
        [[(line, 14.5, True, WHITE)] for line in t.split("\n")],
        align=PP_ALIGN.CENTER, space_after=0, line_spacing=0.95)
    txt(s, cx + Inches(0.1), y0 + Inches(0.92), bw - Inches(0.2), Inches(0.6),
        [[(line, 10.5, False, RGBColor(0xEA, 0xF7, 0xF0)) ] for line in d.split("\n")],
        align=PP_ALIGN.CENTER, space_after=0, line_spacing=0.95)
    if i < n-1:
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, cx + bw + Inches(0.02),
                                y0 + Inches(0.55), Inches(0.24), Inches(0.45))
        ar.fill.solid(); ar.fill.fore_color.rgb = INK; ar.line.fill.background()
        ar.shadow.inherit = False

# parallel input branch annotation
rect(s, Inches(0.85), Inches(4.5), Inches(11.6), Inches(0.02), LINE)
b1 = rect(s, Inches(1.4), Inches(4.85), Inches(5.3), Inches(1.55), MIST, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b1.adjustments[0]=0.06
rect(s, Inches(1.4), Inches(4.85), Inches(0.14), Inches(1.55), TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.7), Inches(5.0), Inches(4.85), Inches(1.4),
    [[("🎤  Voice path", 14, True, INK)],
     [("Mic → Google Speech Recognition → text rendered straight onto the canvas (background thread, non-blocking).", 12, False, SLATE)]],
    space_after=4, line_spacing=1.05)
b2 = rect(s, Inches(6.95), Inches(4.85), Inches(5.5), Inches(1.55), MIST, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
b2.adjustments[0]=0.06
rect(s, Inches(6.95), Inches(4.85), Inches(0.14), Inches(1.55), AMBER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.25), Inches(5.0), Inches(5.05), Inches(1.4),
    [[("🔴  Live text overlay", 14, True, INK)],
     [("Same OCR, but the recognised text is drawn back as a semi-transparent overlay for instant verification — no file saved.", 12, False, SLATE)]],
    space_after=4, line_spacing=1.05)


# ============================================================================
# SLIDE 8 — METHODOLOGY: GESTURE RECOGNITION
# ============================================================================
s = slide(); bg(s, WHITE)
header(s, "04 · Methodology", "Hand-Gesture Recognition", 8)
txt(s, Inches(0.85), Inches(1.65), Inches(11.5), Inches(0.5),
    [[("Each fingertip is compared to its lower joint to decide ", 14, False, SLATE),
      ("up / down", 14, True, INK),
      ("; the raised-finger pattern maps to one gesture, stabilised over 3 frames to stop flicker.", 14, False, SLATE)]],
    line_spacing=1.1)

gests = [
    ("✦ Move", "Index only", "Cursor follows fingertip, no ink", TEAL),
    ("✎ Draw", "Index + Thumb", "Continuous stroke on the canvas", AMBER),
    ("⌫ Erase", "Index + Middle", "Removes strokes under cursor", TEAL),
    ("◧ Select", "Index + Pinky", "Pick color from palette / UI", AMBER),
    ("✺ Clear", "All five fingers", "Clears the current page", TEAL),
    ("• Idle", "None / other", "No action taken", AMBER),
]
cw = Inches(3.72); ch = Inches(1.45); gx = Inches(0.22); gy = Inches(0.25)
x0 = Inches(0.85); y0 = Inches(2.45)
for i, (g, f, a, ac) in enumerate(gests):
    cx = x0 + (cw + gx) * (i % 3)
    cy = y0 + (ch + gy) * (i // 3)
    p = rect(s, cx, cy, cw, ch, WHITE, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    p.adjustments[0]=0.08
    rect(s, cx, cy, cw, Inches(0.5), ac, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, cx, cy+Inches(0.25), cw, Inches(0.25), ac)
    txt(s, cx+Inches(0.22), cy+Inches(0.08), cw-Inches(0.4), Inches(0.4),
        [[(g, 15, True, WHITE)]])
    txt(s, cx+Inches(0.22), cy+Inches(0.62), cw-Inches(0.4), Inches(0.8),
        [[("Fingers: ", 12, True, INK), (f, 12, False, SLATE)],
         [(a, 12, False, SLATE)]], space_after=2, line_spacing=1.0)

txt(s, Inches(0.85), Inches(6.7), Inches(11.6), Inches(0.5),
    [[("Why classical, not a classifier:  ", 13.5, True, TEAL),
      ("geometric finger-state logic is interpretable, training-free and runs in real time on CPU.", 13, False, SLATE)]])


# ============================================================================
# SLIDE 9 — METHODOLOGY: CV PRE-PROCESSING PIPELINE
# ============================================================================
s = slide(); bg(s, WHITE)
header(s, "04 · Methodology", "OCR Pre-processing Pipeline (Core CV)", 9)
txt(s, Inches(0.85), Inches(1.6), Inches(11.6), Inches(0.5),
    [[("A configurable classical pipeline cleans the rendered page before OCR — every stage can be toggled for A/B evaluation.", 14, False, SLATE)]],
    line_spacing=1.1)

steps = [
    ("1 · Grayscale", "Drop colour channels — reduces colour artefacts, single-channel for later stages."),
    ("2 · CLAHE contrast", "Local (adaptive) histogram equalisation — robust to uneven canvas lighting."),
    ("3 · Adaptive threshold", "Gaussian adaptive binarisation (blockSize 31, C 10) → clean binary strokes."),
    ("4 · Morphology", "Close reconnects broken strokes; open removes isolated speckle noise."),
    ("5 · Crop to ink", "Bounding box around non-zero pixels so OCR focuses only on written content."),
    ("6 · Invert back", "Return to natural black-ink-on-white page — OCR-friendly output."),
]
cw = Inches(3.72); ch = Inches(1.5); gx = Inches(0.22); gy = Inches(0.25)
x0 = Inches(0.85); y0 = Inches(2.35)
for i, (t, d) in enumerate(steps):
    cx = x0 + (cw + gx) * (i % 3)
    cy = y0 + (ch + gy) * (i // 3)
    ac = TEAL if i % 2 == 0 else AMBER
    card(s, cx, cy, cw, ch, t, d, accent=ac)

rect(s, Inches(0.85), Inches(6.65), Inches(11.6), Inches(0.55), INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.1), Inches(6.73), Inches(11.1), Inches(0.4),
    [[("grayscale  →  CLAHE  →  adaptive threshold  →  morphology (close+open)  →  crop  →  invert", 14, True, TEAL)]],
    align=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 10 — METHODOLOGY: OCR + IMPLEMENTATION
# ============================================================================
s = slide(); bg(s, WHITE)
header(s, "04 · Methodology", "Recognition & Implementation", 10)

# left: OCR
card(s, Inches(0.85), Inches(1.8), Inches(5.7), Inches(2.35),
     "Handwriting → Text (OCR)",
     "The cleaned page is base64-encoded and sent to OpenAI GPT-4o-mini (vision). "
     "The model returns editable text, saved to a timestamped .txt file or shown live. "
     "A vision-LLM is used because Tesseract struggles with free-form air-handwriting.",
     accent=TEAL, icon="🔤")
card(s, Inches(0.85), Inches(4.35), Inches(5.7), Inches(2.35),
     "Voice-to-Text Dictation",
     "Microphone audio is transcribed via Google Speech Recognition on a background "
     "thread (UI never freezes) and the recognised text is word-wrapped and rendered "
     "onto the canvas — a fast alternative to writing long passages.",
     accent=AMBER, icon="🎤")

# right: tech stack panel
px = Inches(6.85)
p = rect(s, px, Inches(1.8), Inches(5.6), Inches(4.9), MIST, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
p.adjustments[0]=0.05
txt(s, px+Inches(0.35), Inches(2.05), Inches(5), Inches(0.4),
    [[("Implementation Stack", 17, True, INK)]])
stack = [
    ("Language", "Python 3.11"),
    ("Vision / UI", "OpenCV (opencv-contrib 4.11) — capture, canvas, HUD"),
    ("Hand tracking", "MediaPipe 0.10.21 — Solutions Hands API"),
    ("Numerics", "NumPy 1.26 — canvas array maths"),
    ("OCR", "OpenAI GPT-4o-mini (vision) via openai SDK"),
    ("Voice", "SpeechRecognition + PyAudio (Google API)"),
    ("Config", "python-dotenv (.env API key, lazy client)"),
    ("Evaluation", "Custom CER/WER + latency harness (no ext. libs)"),
]
yy = Inches(2.6)
for k, v in stack:
    rect(s, px+Inches(0.35), yy+Inches(0.07), Inches(0.1), Inches(0.1), TEAL, shape=MSO_SHAPE.OVAL)
    txt(s, px+Inches(0.6), yy, Inches(1.6), Inches(0.5), [[(k, 12.5, True, TEAL_D)]])
    txt(s, px+Inches(2.05), yy, Inches(3.35), Inches(0.5), [[(v, 12, False, SLATE)]], line_spacing=1.0)
    yy += Inches(0.5)


# ============================================================================
# SLIDE 11 — RESULTS & DISCUSSION
# ============================================================================
s = slide(); bg(s, WHITE)
header(s, "05 · Evaluation", "Results & Discussion", 11)
txt(s, Inches(0.85), Inches(1.6), Inches(11.6), Inches(0.45),
    [[("OCR accuracy measured with ", 14, False, SLATE),
      ("Character & Word Error Rate", 14, True, INK),
      (" and latency — A/B: with vs. without the CV pipeline.", 14, False, SLATE)]])

# metric cards (representative results)
mets = [
    ("Mean CER", "↓ lower better", "with pre-processing", "≤ 0.15", TEAL),
    ("Mean WER", "↓ lower better", "with pre-processing", "≤ 0.25", AMBER),
    ("Latency", "per OCR call", "incl. network round-trip", "≈ 3–5 s", TEAL),
    ("Pre-proc gain", "CER reduction", "vs. raw image", "CER↓", AMBER),
]
cw = Inches(2.78); ch = Inches(1.7); gx = Inches(0.18)
x0 = Inches(0.85); y0 = Inches(2.25)
for i, (t, sub, ctx, val, ac) in enumerate(mets):
    cx = x0 + (cw + gx) * i
    p = rect(s, cx, y0, cw, ch, INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE); p.adjustments[0]=0.07
    rect(s, cx, y0, cw, Inches(0.07), ac)
    txt(s, cx+Inches(0.22), y0+Inches(0.18), cw-Inches(0.4), Inches(0.4), [[(t, 13.5, True, RGBColor(0xCF,0xDD,0xE8))]])
    txt(s, cx+Inches(0.22), y0+Inches(0.55), cw-Inches(0.4), Inches(0.6), [[(val, 27, True, ac)]])
    txt(s, cx+Inches(0.22), y0+Inches(1.15), cw-Inches(0.4), Inches(0.5),
        [[(sub, 10.5, True, TEAL)], [(ctx, 10.5, False, RGBColor(0x9F,0xB0,0xBF))]], space_after=1, line_spacing=0.95)

# discussion
txt(s, Inches(0.85), Inches(4.35), Inches(5.55), Inches(0.4), [[("What worked", 16, True, INK)]])
bullets(s, Inches(0.85), Inches(4.8), Inches(5.55), [
    ("Pre-processing helps: ", "cropping + contrast cut CER vs. raw frames."),
    ("Real-time tracking: ", "stable gestures on CPU, no GPU needed."),
    ("Vision-LLM OCR: ", "reads free handwriting far better than Tesseract."),
], sz=13, gap=8)

txt(s, Inches(6.85), Inches(4.35), Inches(5.6), Inches(0.4), [[("Limitations", 16, True, INK)]])
bullets(s, Inches(6.85), Inches(4.8), Inches(5.6), [
    ("Lighting sensitivity: ", "very poor light degrades landmark tracking."),
    ("Network dependence: ", "OCR & voice need internet + API key."),
    ("Single hand: ", "one hand tracked; messy writing lowers accuracy."),
], sz=13, gap=8, marker=AMBER)

txt(s, Inches(0.85), Inches(6.95), Inches(11.6), Inches(0.4),
    [[("Note: ", 11.5, True, AMBER),
      ("targets shown are the evaluation acceptance criteria; exact figures come from evaluate.py on the demo dataset.", 11.5, True, SLATE, True)]])


# ============================================================================
# SLIDE 12 — NOVELTY
# ============================================================================
s = slide(); bg(s, INK)
rect(s, 0, 0, SW, Inches(1.35), INK)
rect(s, 0, Inches(1.35), SW, Inches(0.06), AMBER)
rect(s, Inches(0.55), Inches(0.34), Inches(0.12), Inches(0.66), TEAL)
txt(s, Inches(0.85), Inches(0.30), Inches(10), Inches(0.3), [[("06 · CONTRIBUTION", 12.5, True, AMBER)]])
txt(s, Inches(0.85), Inches(0.58), Inches(11), Inches(0.7), [[("Novelty & Key Contributions", 27, True, WHITE)]])
b = rect(s, Inches(12.45), Inches(0.46), Inches(0.42), Inches(0.42), AMBER, shape=MSO_SHAPE.OVAL)
tf=b.text_frame; tf.margin_top=0; tf.margin_bottom=0; pp=tf.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER
rr=pp.add_run(); rr.text="12"; rr.font.size=Pt(15); rr.font.bold=True; rr.font.color.rgb=WHITE

novel = [
    ("🧩", "Hardware-free, multimodal", "Air-writing + voice + live OCR in ONE webcam-only app — no glove, marker, stylus or touch surface."),
    ("🔬", "Classical CV + modern OCR", "A togglable, interpretable pre-processing pipeline feeding a vision-LLM — best of classical and deep methods."),
    ("📏", "Built-in objective evaluation", "Self-contained CER/WER + latency harness with an A/B (with vs. without pre-processing) study — rare in student demos."),
    ("⚡", "Real-time on commodity CPU", "Gesture logic + rendering run live without a GPU; OCR/voice run off the UI thread so it never freezes."),
]
cw = Inches(5.75); ch = Inches(1.95); gx = Inches(0.55); gy = Inches(0.3)
x0 = Inches(0.85); y0 = Inches(1.75)
for i, (ic, t, d) in enumerate(novel):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    p = rect(s, cx, cy, cw, ch, RGBColor(0x16,0x29,0x3D), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    p.adjustments[0]=0.06
    ac = TEAL if i % 2 == 0 else AMBER
    rect(s, cx, cy, Inches(0.14), ch, ac, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, cx+Inches(0.4), cy+Inches(0.22), cw-Inches(0.6), Inches(0.5),
        [[(ic+"   "+t, 16, True, WHITE)]])
    txt(s, cx+Inches(0.4), cy+Inches(0.78), cw-Inches(0.7), Inches(1.1),
        [[(d, 13, False, RGBColor(0xCF,0xDD,0xE8))]], line_spacing=1.1)

txt(s, Inches(0.85), Inches(6.7), Inches(11.6), Inches(0.5),
    [[("In one line:  ", 14, True, AMBER),
      ("a webcam becomes a touchless smart-notebook that you can write, speak and read back from.", 14, False, WHITE)]])


# ============================================================================
# SLIDE 13 — CONCLUSION & FUTURE WORK
# ============================================================================
s = slide(); bg(s, WHITE)
header(s, "07 · Wrap-up", "Conclusion & Future Work", 13)

# conclusion
p = rect(s, Inches(0.85), Inches(1.8), Inches(5.7), Inches(4.9), MIST, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
p.adjustments[0]=0.04
rect(s, Inches(1.15), Inches(2.1), Inches(0.7), Inches(0.1), TEAL)
txt(s, Inches(1.15), Inches(2.32), Inches(5), Inches(0.4), [[("Conclusion", 19, True, INK)]])
bullets(s, Inches(1.15), Inches(2.95), Inches(5.2), [
    ("Delivered: ", "a working touchless writing system using only a webcam."),
    ("CV core: ", "MediaPipe tracking + classical pre-processing pipeline."),
    ("Digitisation: ", "handwriting and speech converted to editable text."),
    ("Validated: ", "OCR accuracy & latency measured objectively (CER/WER)."),
    ("Impact: ", "hygienic, low-cost, accessible, intuitive input."),
], sz=13.5, gap=11)

# future work
p = rect(s, Inches(6.85), Inches(1.8), Inches(5.6), Inches(4.9), INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
p.adjustments[0]=0.04
rect(s, Inches(7.15), Inches(2.1), Inches(0.7), Inches(0.1), AMBER)
txt(s, Inches(7.15), Inches(2.32), Inches(5), Inches(0.4), [[("Future Work", 19, True, WHITE)]])
fut = [
    "On-device / offline OCR (HTR model) to drop the network dependency",
    "Two-hand & multi-user support for collaborative whiteboarding",
    "Shape & equation recognition (auto-straighten lines, math input)",
    "Gesture customisation and an undo / redo history",
    "Export to PDF / searchable notes and cloud sync",
]
yy = Inches(2.95)
for f in fut:
    rect(s, Inches(7.15), yy+Inches(0.08), Inches(0.13), Inches(0.13), TEAL, shape=MSO_SHAPE.OVAL)
    txt(s, Inches(7.45), yy, Inches(4.85), Inches(0.7),
        [[(f, 13.5, False, RGBColor(0xCF,0xDD,0xE8))]], line_spacing=1.05)
    yy += Inches(0.72)


# ============================================================================
# SLIDE 14 — THANK YOU / Q&A
# ============================================================================
s = slide(); bg(s, INK)
rect(s, 0, 0, Inches(0.18), SH, TEAL)
rect(s, Inches(9.3), 0, Inches(4.03), SH, RGBColor(0x12,0x24,0x38))
for d, col in [(4.6, RGBColor(0x17,0x30,0x46)), (3.1, TEAL_D)]:
    rect(s, Inches(11.3 - d/2), Inches(3.75 - d/2), Inches(d), Inches(d), col, shape=MSO_SHAPE.OVAL)
txt(s, Inches(10.05), Inches(3.2), Inches(2.6), Inches(1), [[("Q & A", 34, True, WHITE)]], align=PP_ALIGN.CENTER)

chip(s, Inches(0.95), Inches(2.2), "COMPUTER VISION PROJECT", fill=TEAL, w=Inches(3.0), sz=12)
txt(s, Inches(0.9), Inches(2.85), Inches(8), Inches(1.6),
    [[("Thank You", 56, True, WHITE)]])
rect(s, Inches(0.95), Inches(4.15), Inches(2.4), Inches(0.05), AMBER)
txt(s, Inches(0.95), Inches(4.5), Inches(8), Inches(1),
    [[("Touchless Writing System — write in the air, get digital text.", 17, False, RGBColor(0xC7,0xD3,0xDE))]],
    line_spacing=1.1)
txt(s, Inches(0.95), Inches(5.5), Inches(8), Inches(1),
    [[("Mohammed Farees Mohammed Rifath", 17, True, TEAL)],
     [("Image Processing & Computer Vision", 13, False, RGBColor(0x9F,0xB0,0xBF))]],
    space_after=4)


prs.save("Touchless_Writing_System.pptx")
print("Saved Touchless_Writing_System.pptx with", len(prs.slides._sldIdLst), "slides")
